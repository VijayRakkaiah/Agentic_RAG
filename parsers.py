import fitz
import json
import pandas as pd
from pptx import Presentation
from PIL import Image
import pytesseract

# ---------- PDF Parser ----------
def pdf_parser(file):
    doc = fitz.open(file)

    text = ""

    for page in doc:
        text = text + page.get_text()

    return {
        "text": text,
        "metadata": {
            "file_name": file.name
        }
    }

# ---------- CSV Parser ----------
def csv_parser(filepath):

    df = pd.read_csv(filepath)

    return {
        "text": df.to_string(index=False),
        "metadata": {
            "file_name": filepath.name
        }
    }

# ---------- Image Parser ----------
def image_parser(filepath):

    img = Image.open(filepath)
    text = pytesseract.image_to_string(image=img)

    return {
        "text": text,
        "metadata": {
            "file_name": filepath.name
        }
    }

# ---------- JSON Parser ----------
def json_parser(filepath):

    with open(filepath, 'r', encoding='utf-8') as f:
        data = json.load(f)

    return {
        "text": json.dumps(data, indent=2),
        "metadata": {
            "file_name":filepath.name
        }
    }

# ---------- PPT Parser ----------
def ppt_parser(filepath):

    prs = Presentation(filepath)

    text = ""

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"

    return {
        "text": text,
        "metadata": {
            "file_name": filepath.name
        }
    }